#!/usr/bin/env python3
"""
Build the Ironcrest x Paradox work-test deck.
Clean, minimal, boardroom-professional. Deep slate + muted bronze accent.
16:9. One idea per slide. Tables where the model is shown. Speaker notes everywhere.
Numbers per v2 reconciled working doc. No em-dashes (use " - "). No hype words.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- Theme ----------
SLATE      = RGBColor(0x1C, 0x25, 0x30)   # deep slate (primary dark / titles bg)
SLATE_SOFT = RGBColor(0x2A, 0x36, 0x44)   # softer slate panel
INK        = RGBColor(0x21, 0x2A, 0x33)   # body text on light
MUTE       = RGBColor(0x5C, 0x68, 0x73)   # muted grey text
BRONZE     = RGBColor(0xB0, 0x86, 0x4A)   # muted gold/bronze accent
BRONZE_DK  = RGBColor(0x8C, 0x68, 0x36)
PAPER      = RGBColor(0xF6, 0xF4, 0xEF)   # warm off-white background
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LINE       = RGBColor(0xD9, 0xD3, 0xC7)   # hairline on paper
GREEN      = RGBColor(0x3F, 0x6B, 0x4F)   # restrained "go" tone
TBL_HDR    = SLATE
TBL_ALT    = RGBColor(0xEC, 0xE8, 0xDF)

FONT_H = "Georgia"        # serif for titles - the restrained "grand strategy" nod
FONT_B = "Calibri"        # clean sans for body

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ---------- helpers ----------
def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color, line_color=None, line_w=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = line_w or Pt(0.75)
    sh.shadow.inherit = False
    return sh


def txt(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP):
    """lines: list of dicts {text, size, color, bold, font, align, space_after, space_before, italic}"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(ln.get("space_after", 4))
        p.space_before = Pt(ln.get("space_before", 0))
        if ln.get("line_spacing"):
            p.line_spacing = ln["line_spacing"]
        r = p.add_run()
        r.text = ln["text"]
        r.font.name = ln.get("font", FONT_B)
        r.font.size = Pt(ln.get("size", 18))
        r.font.bold = ln.get("bold", False)
        r.font.italic = ln.get("italic", False)
        r.font.color.rgb = ln.get("color", INK)
    return tb


def bullets(slide, x, y, w, h, items, size=18, color=INK, gap=10, marker=True, lead_color=BRONZE):
    """items: list of (head, body) or plain strings. Bronze square marker."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        if marker:
            m = p.add_run()
            m.text = "▪  "   # small black square
            m.font.name = FONT_B
            m.font.size = Pt(size)
            m.font.color.rgb = lead_color
        if isinstance(it, tuple):
            head, body = it
            rh = p.add_run()
            rh.text = head
            rh.font.name = FONT_B
            rh.font.size = Pt(size)
            rh.font.bold = True
            rh.font.color.rgb = color
            if body:
                rb = p.add_run()
                rb.text = " " + body
                rb.font.name = FONT_B
                rb.font.size = Pt(size)
                rb.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = it
            r.font.name = FONT_B
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb


def content_header(slide, kicker, title, title_size=32):
    """Standard content-slide header on paper bg: bronze kicker + serif title + hairline."""
    bg(slide, PAPER)
    # left accent bar
    rect(slide, Inches(0), Inches(0), Inches(0.18), EMU_H, BRONZE)
    txt(slide, Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.4),
        [{"text": kicker.upper(), "size": 13, "color": BRONZE_DK, "bold": True,
          "font": FONT_B}])
    txt(slide, Inches(0.7), Inches(0.86), Inches(11.9), Inches(0.9),
        [{"text": title, "size": title_size, "color": SLATE, "bold": True,
          "font": FONT_H}])
    rect(slide, Inches(0.72), Inches(1.72), Inches(11.9), Pt(1.5), LINE)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def table(slide, x, y, w, rows, col_widths, header_bold=True, font_size=14,
          header_size=14, row_h=Inches(0.46), align=None):
    """rows: list of lists of strings. First row = header. col_widths in Inches fractions summing to w."""
    nrows = len(rows)
    ncols = len(rows[0])
    gtbl = slide.shapes.add_table(nrows, ncols, x, y, w, row_h * nrows).table
    # disable banding style default; we color manually
    gtbl.first_row = False
    gtbl.horz_banding = False
    # column widths
    for ci, cw in enumerate(col_widths):
        gtbl.columns[ci].width = cw
    for ri in range(nrows):
        gtbl.rows[ri].height = row_h
        for ci in range(ncols):
            cell = gtbl.cell(ri, ci)
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # fill
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TBL_HDR
            elif ri % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TBL_ALT
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            if align and align[ci] is not None:
                p.alignment = align[ci]
            else:
                p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            r = p.add_run()
            r.text = rows[ri][ci]
            r.font.name = FONT_B
            if ri == 0:
                r.font.size = Pt(header_size)
                r.font.bold = True
                r.font.color.rgb = WHITE
            else:
                r.font.size = Pt(font_size)
                r.font.color.rgb = INK
                if ci == 0:
                    r.font.bold = True
    return gtbl


def page_num(slide, n):
    txt(slide, Inches(12.4), Inches(7.02), Inches(0.7), Inches(0.35),
        [{"text": str(n), "size": 10, "color": MUTE, "align": PP_ALIGN.RIGHT, "font": FONT_B}])


# =========================================================
# SLIDE 1 - TITLE
# =========================================================
s = add_slide()
bg(s, SLATE)
# bronze rule top
rect(s, Inches(0.9), Inches(2.35), Inches(2.2), Pt(3), BRONZE)
txt(s, Inches(0.9), Inches(1.55), Inches(11.5), Inches(0.5),
    [{"text": "PARADOX INTERACTIVE  -  BUSINESS OWNER, NEW BUSINESS  -  WORK TEST",
      "size": 14, "color": BRONZE, "bold": True, "font": FONT_B}])
txt(s, Inches(0.86), Inches(2.6), Inches(11.6), Inches(1.8),
    [{"text": "Ironcrest x Paradox", "size": 54, "color": WHITE, "bold": True, "font": FONT_H},
     {"text": "Sign or pass?", "size": 40, "color": BRONZE, "bold": False, "font": FONT_H,
      "space_before": 4}])
txt(s, Inches(0.9), Inches(4.75), Inches(11.5), Inches(0.5),
    [{"text": "Business case + go / no-go", "size": 20, "color": RGBColor(0xC9,0xCE,0xD4),
      "font": FONT_B}])
# footer line
rect(s, Inches(0.9), Inches(6.35), Inches(11.5), Pt(1), SLATE_SOFT)
txt(s, Inches(0.9), Inches(6.5), Inches(8), Inches(0.4),
    [{"text": "[Presenter name]", "size": 15, "color": WHITE, "font": FONT_B}])
txt(s, Inches(8.0), Inches(6.5), Inches(4.4), Inches(0.4),
    [{"text": "22 June 2026", "size": 15, "color": RGBColor(0xA9,0xB0,0xB8),
      "align": PP_ALIGN.RIGHT, "font": FONT_B}])
notes(s, "Open by naming the job: a Business Owner decides whether to sign new titles. "
         "This is that decision, made out loud. The panel said they care about reasoning over numbers, "
         "so I will lead with the call and the thinking, then let the model back it up. "
         "Replace the placeholder with your name.")

# =========================================================
# SLIDE 2 - THE CALL, UP FRONT
# =========================================================
s = add_slide()
bg(s, SLATE)
rect(s, Inches(0), Inches(0), Inches(0.18), EMU_H, BRONZE)
txt(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5),
    [{"text": "THE CALL, UP FRONT", "size": 15, "color": BRONZE, "bold": True, "font": FONT_B}])
txt(s, Inches(0.86), Inches(2.35), Inches(11.6), Inches(1.6),
    [{"text": "Recommendation: Conditional GO", "size": 46, "color": WHITE, "bold": True,
      "font": FONT_H}])
rect(s, Inches(0.92), Inches(3.75), Inches(3.0), Pt(3), BRONZE)
txt(s, Inches(0.9), Inches(4.2), Inches(11.0), Inches(1.2),
    [{"text": "Sign it - structured so the deal prices the risks in, rather than betting on optimism.",
      "size": 24, "color": RGBColor(0xD7,0xDB,0xE0), "font": FONT_B, "line_spacing": 1.15}])
notes(s, "Lead with the answer - a panel reviewing a Business Owner wants the verdict first, then the reasoning. "
         "The whole case hangs on the word 'conditional': the game is a fit, but the structure is what makes it "
         "a responsible yes. I will spend the rest of the deck earning this line.")

# =========================================================
# SLIDE 3 - HOW I'M THINKING ABOUT THIS (the spine)
# =========================================================
s = add_slide()
content_header(s, "The spine", "How I'm thinking about this")
steps = [
    ("01", "What does Paradox specifically bring",
     "that no other publisher can? If the answer is 'nothing special', you self-publish and walk."),
    ("02", "The honest risk",
     "named out loud - including the risk inside Paradox's own portfolio and current strategy."),
    ("03", "A structure that prices the risk in",
     "rather than hand-waving it - the deal terms do the de-risking."),
    ("04", "Success defined as a funnel, not a number",
     "because the Paradox model earns most of its money in the years after launch."),
]
y = 2.1
for num, head, body in steps:
    rect(s, Inches(0.7), Inches(y), Inches(0.7), Inches(0.95), SLATE)
    txt(s, Inches(0.7), Inches(y+0.22), Inches(0.7), Inches(0.6),
        [{"text": num, "size": 24, "color": BRONZE, "bold": True, "align": PP_ALIGN.CENTER, "font": FONT_H}])
    txt(s, Inches(1.65), Inches(y+0.04), Inches(11.0), Inches(0.95),
        [{"text": head, "size": 19, "color": SLATE, "bold": True, "font": FONT_B, "space_after": 2},
         {"text": body, "size": 15, "color": MUTE, "font": FONT_B, "line_spacing": 1.05}])
    y += 1.18
notes(s, "This is the argument in four moves - state it, then deliver it in order. "
         "The point of showing the spine first: the panel sees a Business Owner with a method, not someone "
         "reacting to a brief. Reasoning over numbers is exactly what this slide signals.")
page_num(s, 3)

# =========================================================
# SLIDE 4 - MARKET, THE OPPORTUNITY
# =========================================================
s = add_slide()
content_header(s, "Market  -  1 of 2", "The opportunity")
bullets(s, Inches(0.7), Inches(2.1), Inches(7.4), Inches(4.5),
    [("Small but deep-pocketed niche.", "Grand strategy is loyal and pays for a decade. Comparables bracket the upside: Manor Lords (accessible) 3M+ units; CK3 (deep) sustains a ~EUR 250 DLC library + subscription."),
     ("A real, open lane.", "'Between Manor Lords and CK3' is a genuine gap - deeper than Manor Lords, more approachable than CK3. Paradox does not own it with a fresh IP."),
     ("Signals are healthy, not hyped.", "95K wishlists, 8,200 Discord, 340K reveal views - all organic, zero paid push. Reads as real demand forming."),
     ("Audience overlaps Paradox almost perfectly.", "CK3 + Victoria 3 players, 25-40, PC-only. The single most important line in the brief.")],
    size=16, gap=14)
# right stat panel
rect(s, Inches(8.5), Inches(2.1), Inches(4.1), Inches(4.4), SLATE)
rect(s, Inches(8.5), Inches(2.1), Inches(4.1), Inches(0.12), BRONZE)
stats = [("95K", "Steam wishlists (organic)"),
         ("8,200", "Discord members (organic)"),
         ("340K", "reveal views, positive"),
         ("25-40", "core age, PC-only")]
yy = 2.45
for big, lbl in stats:
    txt(s, Inches(8.75), Inches(yy), Inches(3.6), Inches(0.55),
        [{"text": big, "size": 30, "color": BRONZE, "bold": True, "font": FONT_H}])
    txt(s, Inches(8.75), Inches(yy+0.55), Inches(3.6), Inches(0.4),
        [{"text": lbl, "size": 13, "color": RGBColor(0xC9,0xCE,0xD4), "font": FONT_B}])
    yy += 1.0
notes(s, "The opportunity in one breath: a real audience, forming organically, sitting in an open lane that maps "
         "onto Paradox's existing base. Push the overlap point hardest - it is what makes this worth more to "
         "Paradox than to any other publisher, which is the bridge into slide 6.")
page_num(s, 4)

# =========================================================
# SLIDE 5 - MARKET, THE RISKS
# =========================================================
s = add_slide()
content_header(s, "Market  -  2 of 2", "The risks, named honestly")
bullets(s, Inches(0.7), Inches(2.05), Inches(11.9), Inches(2.3),
    [("Execution risk on the hard 40%.", "Late-game depth, AI, balance, performance - the part that historically slips and breaks games in this genre. The '60% done' is the comfortable 60%."),
     ("Unproven studio at this scope.", "Ashfeld (480K, 78%) shows competence, but a city-builder is a different systems class than a multi-generational grand strategy sim. No external-milestone track record."),
     ("Timeline optimism.", "6mo / 18mo from a team with limited QA and no publisher discipline likely slips 50-100%.")],
    size=15.5, gap=9)
# THE BIG ONE panel
rect(s, Inches(0.7), Inches(4.55), Inches(11.9), Inches(2.0), RGBColor(0x3A,0x2A,0x22))
rect(s, Inches(0.7), Inches(4.55), Inches(0.14), Inches(2.0), BRONZE)
txt(s, Inches(1.0), Inches(4.7), Inches(11.3), Inches(0.4),
    [{"text": "THE BIG ONE  -  STRATEGIC-FIT TENSION", "size": 14, "color": BRONZE, "bold": True, "font": FONT_B}])
txt(s, Inches(1.0), Inches(5.12), Inches(11.4), Inches(1.4),
    [{"text": "Paradox FY2025: operating profit down 80%, a Bloodlines 2 write-down, and a stated pivot to "
              "'deep strategy and management, mostly in-house'. Ironcrest is dead-centre genre fit but cuts against "
              "the in-house half - an external 15-person studio with protective founders.",
      "size": 15.5, "color": RGBColor(0xE7,0xE2,0xD9), "font": FONT_B, "line_spacing": 1.12, "space_after": 4},
     {"text": "Plus: EU5 portfolio cannibalisation (same wallet) and relationship / design-control friction.",
      "size": 14, "color": RGBColor(0xC3,0xBB,0xAD), "font": FONT_B, "italic": True}])
notes(s, "Naming the risks - especially the one inside Paradox's own house - is the move that earns credibility. "
         "The big one is strategic fit: I am effectively asking them to make an external bet right after they "
         "wrote one down. The deal has to answer 'why does this not become the next Bloodlines 2?' - and slide 12 does.")
page_num(s, 5)

# =========================================================
# SLIDE 6 - WHY PARADOX SPECIFICALLY (core)
# =========================================================
s = add_slide()
content_header(s, "The core of the case", "Why Paradox specifically")
txt(s, Inches(0.7), Inches(1.95), Inches(11.9), Inches(0.5),
    [{"text": "A publisher deal is only worth signing if the publisher adds more than the developer gives up. "
              "Here that gap is unusually large:", "size": 15, "color": MUTE, "font": FONT_B, "italic": True,
      "line_spacing": 1.05}])
cards = [
    ("Owned-audience funnel", "CK3 / Vic3 / EU5 / Stellaris base via launcher, mail, forums. De-risks the hardest step - turning wishlists into players - more than money could."),
    ("DLC live-ops competency", "Nobody runs a base + DLC strategy title for a decade better. The late-game gap stops being a risk and turns a free EA roadmap into a paid post-1.0 flywheel."),
    ("Loc + QA pipelines", "The studio explicitly lacks both. These are exactly the gaps a publisher fills, and Paradox has the machinery."),
]
cw = Inches(3.83)
cx = Inches(0.7)
for head, body in cards:
    rect(s, cx, Inches(2.65), cw, Inches(2.7), WHITE, line_color=LINE, line_w=Pt(1))
    rect(s, cx, Inches(2.65), cw, Inches(0.12), BRONZE)
    txt(s, cx+Inches(0.25), Inches(2.95), cw-Inches(0.5), Inches(0.7),
        [{"text": head, "size": 18, "color": SLATE, "bold": True, "font": FONT_B}])
    txt(s, cx+Inches(0.25), Inches(3.65), cw-Inches(0.5), Inches(1.6),
        [{"text": body, "size": 14, "color": INK, "font": FONT_B, "line_spacing": 1.12}])
    cx += Inches(4.05)
rect(s, Inches(0.7), Inches(5.7), Inches(11.9), Inches(0.85), TBL_ALT)
txt(s, Inches(0.95), Inches(5.85), Inches(11.4), Inches(0.6),
    [{"text": "Mutual leverage cuts both ways: Paradox needs titles for its refocused core slate, so the studio "
              "should not overpay or hand over the IP. There is leverage on both sides.",
      "size": 14, "color": INK, "font": FONT_B, "italic": True, "line_spacing": 1.08}])
notes(s, "This is the heart of the case. The argument is not 'good game, let's publish it' - it is 'this specific "
         "publisher converts this specific game's weaknesses into its own strengths'. Land point 2 hardest: the "
         "late-game-content gap, which terrifies the studio, is literally Paradox's business model. The mutual-leverage "
         "note sets up the deal terms - we are not desperate, and neither are they.")
page_num(s, 6)

# =========================================================
# SLIDE 7 - THE COST SIDE
# =========================================================
s = add_slide()
content_header(s, "The model  -  cost", "What Paradox actually risks")
rows = [
    ["Bucket", "Estimate", "Logic"],
    ["Core dev to 1.0", "EUR 1.2 - 1.6M", "15 people x EUR 600-800K/yr x 1.5-2.5 yr. The hard 40% costs above linear."],
    ["Localisation", "EUR 150 - 250K", "4-5 launch languages (the Paradox pattern); rest added live / DLC-funded."],
    ["External QA (PC-only)", "EUR 200K", "Studio lacks QA. ~18-24 mo PC QA. No console cert cost."],
    ["Marketing", "EUR 750K - 1.0M", "~5-7% of base-case gross. Owned audience keeps paid efficient."],
    ["New money to 1.0", "~EUR 2.4 - 3.3M", "What Paradox puts at risk. Lower end = scoped loc."],
]
t = table(s, Inches(0.7), Inches(2.0), Inches(11.9), rows,
          [Inches(2.7), Inches(2.4), Inches(6.8)], font_size=13.5, header_size=14,
          row_h=Inches(0.6))
# all-in callout
rect(s, Inches(0.7), Inches(5.95), Inches(11.9), Inches(0.78), SLATE)
txt(s, Inches(0.95), Inches(6.08), Inches(11.4), Inches(0.55),
    [{"text": "All-in to full release  ~EUR 4.5 - 5.4M   (or ~EUR 4.4 - 4.7M with loc scoped)   "
              "|   Timeline, realistic: EA ~8-10 mo, 1.0 ~24-30 mo",
      "size": 15, "color": WHITE, "bold": True, "font": FONT_B}])
notes(s, "Costs first, deliberately - a Business Owner sizes the bet before the upside. The two numbers to "
         "anchor: ~EUR 2.4-3.3M of NEW money at risk, ~EUR 4.5-5.4M all-in including the studio's EUR 2.1M sunk. "
         "Flag the localisation scoping decision openly - '8-10 languages at EUR 150K' would be indefensible, so we "
         "recommend the 4-5 launch-language Paradox pattern. The realistic timeline (not the brief's optimistic one) "
         "is what we size the advance against.")
page_num(s, 7)

# =========================================================
# SLIDE 8 - THE REVENUE MODEL
# =========================================================
s = add_slide()
content_header(s, "The model  -  revenue", "Scenarios, not a forecast")
rows = [
    ["Scenario", "Units (2yr)", "Net / unit", "Base rev", "+ DLC", "Total (2yr)"],
    ["Bear", "250K", "EUR 13", "EUR 3.25M", "+30%", "~EUR 4.25M"],
    ["Base", "500K", "EUR 15", "EUR 7.5M", "+50%", "~EUR 11.25M"],
    ["Bull", "1.2M", "EUR 16", "EUR 19.2M", "+60%", "~EUR 30.7M"],
]
table(s, Inches(0.7), Inches(2.05), Inches(11.9), rows,
      [Inches(1.9), Inches(1.9), Inches(1.9), Inches(2.0), Inches(1.7), Inches(2.5)],
      font_size=15, header_size=14, row_h=Inches(0.62))
# footnote panel
rect(s, Inches(0.7), Inches(5.1), Inches(11.9), Inches(1.5), TBL_ALT)
rect(s, Inches(0.7), Inches(5.1), Inches(0.12), Inches(1.5), BRONZE)
txt(s, Inches(0.95), Inches(5.22), Inches(11.4), Inches(1.35),
    [{"text": "Conversion logic:", "size": 14, "color": SLATE, "bold": True, "font": FONT_B, "space_after": 2},
     {"text": "Stated lifetime tiers, not a hidden multiplier - Year 1 = 30-90% of the EA wishlist bank "
              "(roughly month-1 x 2-3). All figures post-Steam 30%, before the dev/publisher split.",
      "size": 13.5, "color": INK, "font": FONT_B, "line_spacing": 1.1, "space_after": 3},
     {"text": "Key caveat: only ~21% of EA titles earn more at 1.0 than in their first 30 EA days "
              "(GameDiscoverCo). So the base case assumes only a MODEST 1.0 beat, not a big one.",
      "size": 13.5, "color": BRONZE_DK, "bold": True, "font": FONT_B, "line_spacing": 1.1}])
notes(s, "Scenarios, not a point forecast - that is the honest way to model a new IP. Walk the base row: 500K units "
         "(pulled down from 600K precisely because of the 21% rule), EUR 15 net, +50% DLC. The caveat is the most "
         "important line on the slide and pre-empts the sharpest panel question: I am NOT banking on a big 1.0 spike, "
         "because the data says most titles don't get one. Remember A5 - this is a 2-year window; Paradox's real model "
         "runs 5-10 years, so even bull understates lifetime value.")
page_num(s, 8)

# =========================================================
# SLIDE 9 - DOES IT PENCIL? (the asymmetry)
# =========================================================
s = add_slide()
content_header(s, "The model  -  does it pencil?", "The asymmetry is the argument")
bullets(s, Inches(0.7), Inches(2.1), Inches(11.9), Inches(2.4),
    [("The bear case is thin on the 2-year P&L alone.", "Paradox's ~50% of EUR 4.25M (~EUR 2.1M) only roughly covers the low end of outlay over two years."),
     ("Downside is protected by structure, not by the forecast.", "(a) the milestone-gated advance caps at-risk capital before the EA gate; (b) the post-2-year DLC tail Paradox reliably produces for an engaged strategy base."),
     ("Base and bull are comfortably profitable.", "Base = Paradox's ~EUR 5.6M share of EUR 11.25M. Every scenario improves sharply in years 3-10.")],
    size=16.5, gap=14)
rect(s, Inches(0.7), Inches(5.25), Inches(11.9), Inches(1.3), GREEN)
txt(s, Inches(0.95), Inches(5.45), Inches(11.4), Inches(0.9),
    [{"text": "Downside protected by structure + the DLC tail. Upside is a decade-long genre franchise.",
      "size": 19, "color": WHITE, "bold": True, "font": FONT_B, "space_after": 3},
     {"text": "That asymmetry - capped downside, uncapped upside - is the financial argument for GO.",
      "size": 15, "color": RGBColor(0xD9,0xE5,0xDD), "font": FONT_B}])
notes(s, "This is the slide that converts the model into a decision. The honest admission - the 2-year bear barely "
         "covers outlay - is a strength, not a weakness: it shows I am not hiding behind the rosy case. The yes does "
         "not rest on the forecast; it rests on the asymmetry. Capped downside (gated advance), uncapped upside "
         "(franchise). That is how a Business Owner should think about a portfolio bet.")
page_num(s, 9)

# =========================================================
# SLIDE 10 - PRICING & STRATEGY MOVES
# =========================================================
s = add_slide()
content_header(s, "Strategy", "Pricing and strategy moves")
bullets(s, Inches(0.7), Inches(2.1), Inches(11.9), Inches(4.4),
    [("Raise 1.0 to EUR 39.99.", "EU5 is EUR 59.99; the depth (40-80h campaigns) and the buyer's willingness to pay support a premium. The brief's EUR 34.99 underprices it."),
     ("Keep EA cheap (EUR 27.99 - 29.99).", "Counter-intuitive but deliberate: in the Paradox model, optimise the base for INSTALL-BASE growth, not ASP - every buyer is a decade-long DLC customer. Lower EA price maximises conversion + review velocity."),
     ("Free content roadmap during EA; paid DLC starts at 1.0.", "Charging for gameplay on an unfinished base sours EA communities. Free updates fill the late-game gap in public; the paid free-plus-DLC flywheel begins at 1.0. The only acceptable EA-phase paid item is an optional, non-gameplay cosmetic supporter pack."),
     ("Subscription is a later lever.", "CK3 added a sub years in, once the DLC library was deep enough to feel intimidating. Not a launch move - part of the long arc."),
     ("Keep 'no live service'.", "Right for the genre. Base + DLC is the proven Paradox shape; do not over-engineer.")],
    size=15.5, gap=11)
notes(s, "The pricing slide is where I show I understand the Paradox model, not just pricing in general. The "
         "non-obvious move is point 2: deliberately UNDER-pricing EA. In a base + DLC model the install base is the "
         "asset, not the launch ASP - so you optimise EA for the number of future DLC customers, not for revenue per "
         "unit. Saying that out loud in the room signals Paradox-native thinking. "
         "Likely panel question - 'do you sell paid DLC during EA?' Answer: no. During EA we ship only FREE content "
         "updates - that is the building-in-public model the grand-strategy audience expects, and it fills the "
         "late-game gap where everyone can see it. Charging for gameplay on top of an unfinished base reliably sours "
         "an EA community, so the paid free-plus-paid expansion flywheel starts at 1.0, not before. The single "
         "exception is an optional, non-gameplay cosmetic / supporter pack (the StarRupture supporter-pack pattern) - "
         "buyers who want to fund development can, and no one who skips it is locked out of any gameplay.")
page_num(s, 10)

# =========================================================
# SLIDE 11 - SUCCESS = A FUNNEL (3 gates)
# =========================================================
s = add_slide()
content_header(s, "Definition of success", "Success is a funnel, measured at three gates")
gates = [
    ("GATE 1  -  6mo / EA", "Did EA validate conversion + stickiness?",
     ["Wishlists 250K+", "EA month-1: 50-75K units", "80%+ reviews", "D30 retention high", "Refund rate <8%"]),
    ("GATE 2  -  Launch / 1.0", "Did it hold up and start the live model?",
     ["500K+ cumulative units", "Review uplift vs EA", "First DLC attach 20%+", "Advance near recouped"]),
    ("GATE 3  -  2 years", "Is this a franchise?",
     ["600K-1M+ lifetime units", "DLC 30%+ of revenue, climbing", "Returning DAU on DLC drops", "Is-it-a-franchise call"]),
]
cw = Inches(3.83)
cx = Inches(0.7)
for title_g, q, items in gates:
    rect(s, cx, Inches(2.05), cw, Inches(4.5), WHITE, line_color=LINE, line_w=Pt(1))
    rect(s, cx, Inches(2.05), cw, Inches(0.72), SLATE)
    txt(s, cx+Inches(0.2), Inches(2.2), cw-Inches(0.4), Inches(0.5),
        [{"text": title_g, "size": 15, "color": BRONZE, "bold": True, "font": FONT_B}])
    txt(s, cx+Inches(0.2), Inches(2.92), cw-Inches(0.4), Inches(0.85),
        [{"text": q, "size": 13.5, "color": MUTE, "italic": True, "font": FONT_B, "line_spacing": 1.05}])
    bullets(s, cx+Inches(0.2), Inches(3.8), cw-Inches(0.4), Inches(2.5),
            items, size=13.5, gap=9)
    cx += Inches(4.05)
notes(s, "Success is not one number - it is whether each stage validates before the next tranche of spend. This "
         "directly mirrors the milestone-gated advance: the funnel and the cash structure are the same shape. The "
         "2-year question is the real one - not 'did it sell' but 'do we commit to a 5-10 year franchise'. That is "
         "the prize.")
page_num(s, 11)

# =========================================================
# SLIDE 12 - RECOMMENDATION + 6 CONDITIONS
# =========================================================
s = add_slide()
content_header(s, "The decision", "Conditional GO  -  the six conditions")
conds = [
    ("Milestone-gated advance EUR 1.0-1.5M.", "Pre-EA tranche bridges runway; EA is the first hard go/no-go gate; post-EA top-up tied to EA performance."),
    ("Dev share 50-55%, concurrent recoup, studio keeps IP.", "Paradox takes a publishing licence + DLC/sequel rights. The franchise is the prize - don't fight over ownership."),
    ("Pricing per the strategy slide.", "Premium EUR 39.99 1.0; install-base-optimised EUR 27.99-29.99 EA."),
    ("Localisation scoped to 4-5 launch languages.", "The Paradox pattern; rest added live / DLC-funded."),
    ("Free EA roadmap; paid DLC flywheel from 1.0.", "Free updates fill the late-game gap during EA; the paid expansion model starts at 1.0. De-risks launch and seeds the live model."),
    ("Design-authority split agreed up front.", "Because Task 2 is coming, and the relationship is the asset behind the 10-year flywheel."),
]
col_w = Inches(5.85)
xs = [Inches(0.7), Inches(6.75)]
for i, (head, body) in enumerate(conds):
    col = i // 3
    row = i % 3
    x = xs[col]
    y = Inches(2.05 + row*1.55)
    rect(s, x, y, Inches(0.55), Inches(1.35), SLATE)
    txt(s, x, y+Inches(0.4), Inches(0.55), Inches(0.6),
        [{"text": str(i+1), "size": 22, "color": BRONZE, "bold": True, "align": PP_ALIGN.CENTER, "font": FONT_H}])
    txt(s, x+Inches(0.72), y+Inches(0.05), col_w-Inches(0.8), Inches(1.4),
        [{"text": head, "size": 14.5, "color": SLATE, "bold": True, "font": FONT_B, "space_after": 2, "line_spacing": 1.0},
         {"text": body, "size": 12.5, "color": MUTE, "font": FONT_B, "line_spacing": 1.02}])
notes(s, "The conditions ARE the answer to 'why won't this be the next write-down?'. Each one prices a named risk: "
         "the gated advance caps capital before the funnel proves out; keeping IP with the studio reflects how 96%+ "
         "of deals work and keeps the relationship healthy; the design-authority split pre-empts Task 2. If the studio "
         "refuses EA entirely, the thesis bends toward the closed-dev, protective-founder pattern Paradox just wrote "
         "down - at which point it becomes a leadership call, not a deal I push on optimism.")
page_num(s, 12)

# =========================================================
# SLIDE 13 - SECTION DIVIDER, TASK 2
# =========================================================
s = add_slide()
bg(s, SLATE)
rect(s, Inches(0.9), Inches(2.9), Inches(2.2), Pt(3), BRONZE)
txt(s, Inches(0.9), Inches(2.15), Inches(11.5), Inches(0.5),
    [{"text": "TASK 2", "size": 16, "color": BRONZE, "bold": True, "font": FONT_B}])
txt(s, Inches(0.86), Inches(3.2), Inches(11.6), Inches(1.2),
    [{"text": "The hard conversation", "size": 46, "color": WHITE, "bold": True, "font": FONT_H}])
txt(s, Inches(0.9), Inches(4.55), Inches(11.2), Inches(1.0),
    [{"text": "Six months in, the studio wants 12 more months of closed development before any "
              "player-facing release.", "size": 19, "color": RGBColor(0xC9,0xCE,0xD4), "font": FONT_B,
      "line_spacing": 1.15}])
txt(s, Inches(0.9), Inches(6.4), Inches(11.2), Inches(0.5),
    [{"text": "Discussed verbally - the next slides are backup / discussion support.",
      "size": 14, "color": RGBColor(0x8A,0x93,0x9C), "italic": True, "font": FONT_B}])
notes(s, "Transition slide. Say plainly that Task 2 is a conversation, not a slide-driven pitch - the next two "
         "slides are there if the discussion needs an anchor. Then put the deck down and talk.")

# =========================================================
# SLIDE 14 - TASK 2, MY RESPONSE
# =========================================================
s = add_slide()
content_header(s, "Task 2", "My response to the studio")
bullets(s, Inches(0.7), Inches(2.05), Inches(11.9), Inches(4.5),
    [("Separate two things they're conflating.", "'More polish' is legitimate. 'No release for 12 months' is a change of commercial strategy, not a schedule slip. Name the distinction calmly."),
     ("Answer the fear, not with authority.", "Under 'we need more time' sits 'EA will tarnish the game'. The fix is NARROWER EA, not LATER EA - ship fewer systems fully polished, rather than everything half-done."),
     ("Bring the cash curve.", "Advance + EA fund the back half jointly. Remove EA and the burden shifts onto a bigger advance on worse terms for them. Show the curve - EA becomes survival + validation, not exposure."),
     ("Use the data.", "Only ~21% of EA titles earn more at 1.0 than in their first 30 EA days (GameDiscoverCo). The commercial peak is usually EA itself - waiting leaves the biggest moment on the table."),
     ("Offer real compromise.", "A 3-month scoped slip to a polished vertical-slice EA - not 12 - with any extension tied to milestone evidence, not vibes.")],
    size=15, gap=10)
notes(s, "The core reframe is 'narrower not later' - it honours their quality fear while protecting the EA gate. "
         "The cash curve is the unlock: they think EA is exposure; show them it's the thing that lets them avoid a "
         "bigger, worse advance. The 21% data point does double duty - it kept us conservative on revenue AND it is "
         "the argument against waiting for a 'perfect 1.0'. Offer the 3-month slip so they get a real concession.")
page_num(s, 14)

# =========================================================
# SLIDE 15 - TASK 2, BALANCE + ESCALATE vs OWN
# =========================================================
s = add_slide()
content_header(s, "Task 2", "Balance, and what I own vs escalate")
# balance band
rect(s, Inches(0.7), Inches(1.95), Inches(11.9), Inches(1.1), TBL_ALT)
rect(s, Inches(0.7), Inches(1.95), Inches(0.12), Inches(1.1), BRONZE)
txt(s, Inches(0.95), Inches(2.08), Inches(11.4), Inches(0.95),
    [{"text": "The relationship is the asset that makes the 10-year flywheel possible. Don't pull rank. "
              "Find the smallest change that protects the commercial logic.",
      "size": 16, "color": INK, "bold": True, "font": FONT_B, "line_spacing": 1.1}])
# two columns
rect(s, Inches(0.7), Inches(3.3), Inches(5.85), Inches(3.25), WHITE, line_color=LINE, line_w=Pt(1))
rect(s, Inches(0.7), Inches(3.3), Inches(5.85), Inches(0.6), GREEN)
txt(s, Inches(0.95), Inches(3.42), Inches(5.4), Inches(0.4),
    [{"text": "I OWN", "size": 16, "color": WHITE, "bold": True, "font": FONT_B}])
bullets(s, Inches(0.95), Inches(4.05), Inches(5.4), Inches(2.4),
    ["EA-scope renegotiation ('narrower not later')",
     "Cash-curve modelling + walking them through it",
     "The relationship and finding the compromise",
     "Success criteria for any agreed slip"],
    size=13.5, gap=10, lead_color=GREEN)

rect(s, Inches(6.75), Inches(3.3), Inches(5.85), Inches(3.25), WHITE, line_color=LINE, line_w=Pt(1))
rect(s, Inches(6.75), Inches(3.3), Inches(5.85), Inches(0.6), SLATE)
txt(s, Inches(7.0), Inches(3.42), Inches(5.4), Inches(0.4),
    [{"text": "I ESCALATE", "size": 16, "color": BRONZE, "bold": True, "font": FONT_B}])
bullets(s, Inches(7.0), Inches(4.05), Inches(5.4), Inches(2.4),
    ["Any change to advance / capital structure",
     "A thesis-breaking timeline change",
     "A genuine EA refusal - starts to look like the Bloodlines 2 risk pattern"],
    size=13.5, gap=10, lead_color=SLATE)
notes(s, "The one-line version: protect the EA gate, protect the relationship, and be honest about which decisions "
         "are mine and which belong to the people who own the capital. The escalate column is the maturity signal - "
         "a Business Owner knows the edge of their authority. A genuine EA refusal isn't mine to override; it goes up, "
         "because that is exactly the risk pattern Paradox just paid for.")
page_num(s, 15)

# =========================================================
# SLIDE 16 - APPENDIX, ASSUMPTIONS & SOURCES
# =========================================================
s = add_slide()
content_header(s, "Appendix", "Key assumptions and sources")
# left: assumptions
txt(s, Inches(0.7), Inches(2.0), Inches(6.0), Inches(0.4),
    [{"text": "STATED ASSUMPTIONS", "size": 14, "color": BRONZE_DK, "bold": True, "font": FONT_B}])
bullets(s, Inches(0.7), Inches(2.5), Inches(6.0), Inches(4.0),
    ["Wishlists ~250K at EA (base case)",
     "Scenario net / unit: EUR 13 / 15 / 16 (bear / base / bull)",
     "Studio burn EUR 600-800K/yr (15 people, Wroclaw)",
     "Timeline slip 50-100% vs the brief",
     "Treated as base + DLC over a 5-10 year horizon"],
    size=14.5, gap=12)
# right: sources
rect(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(4.5), SLATE)
txt(s, Inches(7.3), Inches(2.2), Inches(5.0), Inches(0.4),
    [{"text": "SOURCES", "size": 14, "color": BRONZE, "bold": True, "font": FONT_B}])
bullets(s, Inches(7.3), Inches(2.7), Inches(5.0), Inches(3.6),
    ["GameDiscoverCo (EA graduation + conversion data)",
     "Alinea Analytics (DLC attach estimates)",
     "Paradox FY2025 year-end report + EU5 launch",
     "CK3 + Manor Lords sales / wishlist data",
     "Our internal benchmarks (Rust Racers model)"],
    size=14.5, gap=12, color=RGBColor(0xD7,0xDB,0xE0), lead_color=BRONZE)
notes(s, "Backup slide. Every number on the deck traces to one of these. The point of listing assumptions explicitly: "
         "if a panellist disagrees with one, the call can move - and saying so out loud is the honest, Business-Owner "
         "thing to do. Don't read it; leave it up if the discussion turns to methodology.")
page_num(s, 16)

# ---------- save ----------
out = "/home/assistant/projects/paradox_ironcrest_case/drafts/ironcrest_paradox_deck.pptx"
prs.save(out)
print("Saved:", out)
print("Slides:", len(prs.slides._sldIdLst))
